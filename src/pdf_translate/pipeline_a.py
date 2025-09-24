"""Pipeline A: convert PDF to PPTX via pdf2pptx, replace text, export results."""
from __future__ import annotations

import logging
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import List, TYPE_CHECKING

from pptx import Presentation  # type: ignore
from tqdm import tqdm

from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
from pptx.enum.text import MSO_AUTO_SIZE  # type: ignore


if TYPE_CHECKING:
    from pptx.text.text import TextFrame
from .config import AppConfig
from .translation import TranslationClient
from .utils import ensure_parent, clean_directory

logger = logging.getLogger(__name__)


@dataclass
class ParagraphHandle:
    paragraph: "Paragraph"
    text: str
    text_frame: "TextFrame"
    original_auto_size: MSO_AUTO_SIZE | None


def _has_hangul(text: str) -> bool:
    return any("\uac00" <= ch <= "\ud7a3" for ch in text)


def _sample_preview(texts: List[str], limit: int = 40, max_items: int = 3) -> str:
    sample = []
    for text in texts[:max_items]:
        collapsed = " ".join(text.split())
        if len(collapsed) > limit:
            collapsed = f"{collapsed[:limit]}…"
        sample.append(collapsed)
    return ", ".join(sample)


class PipelineA:
    def __init__(self, config: AppConfig, translator: TranslationClient, work_dir: Path | None = None) -> None:
        self.config = config
        self.translator = translator
        self.work_dir = work_dir or config.working_dir
        self.work_dir.mkdir(parents=True, exist_ok=True)

    def run(self) -> Path:
        source_path = self.config.input_path
        pptx_path = self._prepare_pptx_source(source_path)
        translated_pptx = self._translate_pptx(pptx_path)
        output_path = self._persist_output(translated_pptx)

        if self.config.cleanup_working:
            logger.debug("Cleaning working directory %s", self.work_dir)
            clean_directory(self.work_dir)

        return output_path

    def _prepare_pptx_source(self, path: Path) -> Path:
        suffix = path.suffix.lower()
        if suffix not in {".ppt", ".pptx"}:
            raise ValueError(
                f"Pipeline A expects PPT/PPTX input, got {path}. Use a PPT source or allow the CLI to route this file through Pipeline B."
            )
        target = self.work_dir / path.name
        shutil.copy2(path, target)
        return target

    def _translate_pptx(self, pptx_path: Path) -> Path:
        logger.info("Translating text content within %s", pptx_path)
        presentation = Presentation(str(pptx_path))
        handles = self._collect_paragraphs(presentation)
        texts = [handle.text for handle in handles]
        logger.info("Identified %d paragraphs containing Hangul", len(texts))
        if texts:
            logger.debug("Sample source paragraphs: %s", _sample_preview(texts))
        translations = []
        # for text in tqdm(texts, desc="Translating paragraphs"):
        #     translation = self.translator.translate_text(text)
        #     translations.append(translation)
        translations = self.translator.translate_batch(texts)

        if translations:
            logger.debug("Sample translated paragraphs: %s", _sample_preview(translations))
        for handle, translation in zip(handles, translations):
            self._apply_translation(handle, translation)
        translated_path = self.work_dir / f"{pptx_path.stem}_translated.pptx"
        presentation.save(str(translated_path))
        return translated_path

    def _persist_output(self, translated_pptx: Path) -> Path:
        target = self.config.output_path
        suffix = target.suffix.lower()
        if suffix not in {".ppt", ".pptx"}:
            adjusted = target.with_suffix(".pptx")
            logger.warning(
                "Pipeline A expects PPT/PPTX output. Writing translated deck to %s instead of %s.",
                adjusted,
                target,
            )
            target = adjusted
        ensure_parent(target)
        shutil.copy2(translated_pptx, target)
        logger.info("Translated PPTX saved to %s", target)
        return target

    def _collect_paragraphs(self, presentation: Presentation) -> List[ParagraphHandle]:
        handles: List[ParagraphHandle] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                self._gather_paragraphs_from_shape(shape, handles)
        return handles

    def _gather_paragraphs_from_shape(self, shape, handles: List[ParagraphHandle]) -> None:
        if getattr(shape, "has_text_frame", False):
            self._gather_from_text_frame(shape.text_frame, handles)
        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    text_frame = getattr(cell, "text_frame", None)
                    if text_frame is not None:
                        self._gather_from_text_frame(text_frame, handles)
        shape_type = getattr(shape, "shape_type", None)
        if shape_type == MSO_SHAPE_TYPE.GROUP:
            for nested_shape in shape.shapes:
                self._gather_paragraphs_from_shape(nested_shape, handles)
        if getattr(shape, "has_chart", False):
            self._gather_from_chart(shape.chart, handles)

    def _gather_from_text_frame(self, text_frame, handles: List[ParagraphHandle]) -> None:
        for paragraph in text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            if not _has_hangul(text):
                continue
            auto_size = getattr(text_frame, "auto_size", None)
            handles.append(ParagraphHandle(paragraph=paragraph, text=text, text_frame=text_frame, original_auto_size=auto_size))

    def _gather_from_chart(self, chart, handles: List[ParagraphHandle]) -> None:
        if chart is None:
            return
        if getattr(chart, "has_title", False):
            chart_title = getattr(chart, "chart_title", None)
            if chart_title is not None:
                self._gather_from_text_frame(chart_title.text_frame, handles)
        for axis_name in ("category_axis", "value_axis", "series_axis"):
            axis = getattr(chart, axis_name, None)
            if axis is None:
                continue
            if getattr(axis, "has_title", False):
                axis_title = getattr(axis, "axis_title", None)
                if axis_title is not None:
                    self._gather_from_text_frame(axis_title.text_frame, handles)

    def _ensure_text_frame_bounds(self, handle: ParagraphHandle, translation: str) -> None:
        if len(translation) <= len(handle.text):
            return
        if handle.original_auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
            return
        text_frame = handle.text_frame
        try:
            if hasattr(text_frame, "word_wrap") and text_frame.word_wrap is False:
                text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except (AttributeError, TypeError, ValueError):
            return

    def _apply_translation(self, handle: ParagraphHandle, translation: str) -> None:
        paragraph = handle.paragraph
        if paragraph.runs:
            paragraph.runs[0].text = translation
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            run = paragraph.add_run()
            run.text = translation
        self._ensure_text_frame_bounds(handle, translation)
